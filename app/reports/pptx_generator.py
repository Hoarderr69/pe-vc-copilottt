"""
PPTX Generator — build an editable 10-slide Board Pack with python-pptx.

Consumes the exact same per-slide dicts produced by ``slide_data_builder.build_slides``
(persisted on ``ReportRecord.slides``), so the PPTX, the in-app viewer, and the
reportlab PDF all render from one source of truth.

Design notes
------------
* Charts (EBITDA forward curve, IRR scenario heatmap) are rendered to PNG with
  **matplotlib** and inserted as pictures. The Phase-3 plan named Playwright for
  this; matplotlib is used instead because it needs no headless-browser binary
  download, runs fully offline, and produces a deterministic PNG that embeds
  identically. Everything else on each slide is native PPTX text/tables, so the
  deck stays fully editable in PowerPoint / LibreOffice.
* Every slide carries the spec chrome: 40px navy (#0F2044) header bar, white body,
  and a light-grey footer reading "Page N / 10 · CONFIDENTIAL — BOARD ONLY".
* Until the report is approved a translucent "DRAFT" watermark is stamped across
  every slide; on approval it is removed.
"""

from __future__ import annotations

import io
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Palette (mirrors the viewer / spec)
# ---------------------------------------------------------------------------

NAVY        = RGBColor(0x0F, 0x20, 0x44)
INK         = RGBColor(0x0F, 0x17, 0x2A)
SLATE       = RGBColor(0x47, 0x55, 0x69)
MUTED       = RGBColor(0x94, 0xA3, 0xB8)
FOOTER_BG   = RGBColor(0xF1, 0xF5, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_FILL = RGBColor(0xE2, 0xE8, 0xF0)

RAG = {
    "red":   RGBColor(0xDC, 0x26, 0x26),
    "amber": RGBColor(0xD9, 0x77, 0x06),
    "green": RGBColor(0x16, 0xA3, 0x4A),
    "grey":  RGBColor(0x94, 0xA3, 0xB8),
}
RAG_SOFT = {
    "red":   RGBColor(0xFE, 0xE2, 0xE2),
    "amber": RGBColor(0xFE, 0xF3, 0xC7),
    "green": RGBColor(0xDC, 0xFC, 0xE7),
    "grey":  RGBColor(0xE2, 0xE8, 0xF0),
}

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.55)
FOOTER_H = Inches(0.38)
MARGIN = Inches(0.55)
BODY_TOP = Inches(0.95)


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _txbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    return box, tf


def _run(paragraph, text, *, size=12, bold=False, color=INK, italic=False, font="Calibri"):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ---------------------------------------------------------------------------
# Chart rendering (matplotlib → PNG bytes)
# ---------------------------------------------------------------------------

def _forward_curve_png(data: Dict[str, Any]) -> Optional[bytes]:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception:
        return None

    series = data.get("series") or {}
    actual = series.get("actual") or []
    proj = series.get("projection") or []
    if not actual:
        return None

    ax_actual = [p.get("ebitda") for p in actual if p.get("ebitda") is not None]
    n_act = len(ax_actual)
    x_act = list(range(n_act))

    fig, ax = plt.subplots(figsize=(8.6, 3.7), dpi=150)
    ax.plot(x_act, ax_actual, color="#0F2044", lw=2.2, marker="o", ms=4, label="Actual EBITDA")

    if proj:
        x_proj = list(range(n_act - 1, n_act - 1 + len(proj) + 1))
        p50 = [ax_actual[-1]] + [p.get("p50") for p in proj]
        p10 = [ax_actual[-1]] + [p.get("p10") for p in proj]
        p90 = [ax_actual[-1]] + [p.get("p90") for p in proj]
        ax.plot(x_proj, p50, color="#2563EB", lw=2, ls="--", marker="o", ms=3, label="P50 forecast")
        ax.fill_between(x_proj, p10, p90, color="#2563EB", alpha=0.12, label="P10–P90 band")

    ic_target = data.get("ic_target")
    if ic_target:
        ax.axhline(ic_target, color="#16A34A", lw=1.4, ls=":", label="IC exit target")

    ax.set_ylabel("EBITDA (£/mo)", fontsize=9, color="#475569")
    ax.tick_params(labelsize=8, colors="#475569")
    ax.grid(True, axis="y", color="#E2E8F0", lw=0.7)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    ax.legend(fontsize=8, loc="upper left", frameon=False)
    fig.tight_layout(pad=0.6)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


def _irr_heatmap_png(matrix: Dict[str, Any]) -> Optional[bytes]:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except Exception:
        return None

    grid = matrix.get("grid") or []
    hold_years = matrix.get("hold_years") or []
    if not grid or not hold_years:
        return None

    rag_color = {"green": "#16A34A", "amber": "#D97706", "red": "#DC2626"}

    fig, ax = plt.subplots(figsize=(8.0, 3.9), dpi=150)
    n_rows = len(grid)
    n_cols = len(hold_years)
    # Rows top-to-bottom = highest multiple first
    for ri, row in enumerate(reversed(grid)):
        for ci, cell in enumerate(row.get("cells", [])):
            color = rag_color.get(cell.get("rag", "grey"), "#94A3B8")
            ax.add_patch(plt.Rectangle((ci, ri), 1, 1, facecolor=color, edgecolor="white", lw=2, alpha=0.92))
            label = cell.get("irr_pct", "—")
            if cell.get("ic_cell"):
                label += "\n(IC)"
            ax.text(ci + 0.5, ri + 0.5, label, ha="center", va="center",
                    fontsize=8.5, color="white", fontweight="bold")

    ax.set_xlim(0, n_cols)
    ax.set_ylim(0, n_rows)
    ax.set_xticks([i + 0.5 for i in range(n_cols)])
    ax.set_xticklabels([f"Yr {y}" for y in hold_years], fontsize=8.5, color="#475569")
    ax.set_yticks([i + 0.5 for i in range(n_rows)])
    ax.set_yticklabels([r.get("exit_multiple_label", "") for r in reversed(grid)],
                       fontsize=8.5, color="#475569")
    ax.set_xlabel("Hold period", fontsize=9, color="#475569")
    ax.set_ylabel("Exit multiple", fontsize=9, color="#475569")
    ax.tick_params(length=0)
    for s in ax.spines.values():
        s.set_visible(False)
    fig.tight_layout(pad=0.6)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Slide chrome
# ---------------------------------------------------------------------------

def _add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _chrome(slide, *, company: str, period: str, page: int, total: int,
            section: str, is_draft: bool):
    # Header bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, HEADER_H)  # 1 = rectangle
    _fill(bar, NAVY)
    _, tf = _txbox(slide, MARGIN, 0, SLIDE_W - 2 * MARGIN, HEADER_H, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _run(p, "VCP COPILOT", size=11, bold=True, color=WHITE)
    p2 = tf.add_paragraph()  # right-aligned company/period via separate box
    # Centre company/period
    _, ctf = _txbox(slide, Inches(3), 0, SLIDE_W - Inches(6), HEADER_H, anchor=MSO_ANCHOR.MIDDLE)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(ctf.paragraphs[0], f"{company}  ·  {period}", size=12, color=WHITE)
    # Right CONFIDENTIAL
    _, rtf = _txbox(slide, SLIDE_W - Inches(3) - MARGIN, 0, Inches(3), HEADER_H, anchor=MSO_ANCHOR.MIDDLE)
    rtf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _run(rtf.paragraphs[0], "CONFIDENTIAL", size=9, color=MUTED)

    # Footer bar
    foot = slide.shapes.add_shape(1, 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H)
    _fill(foot, FOOTER_BG)
    _, lf = _txbox(slide, MARGIN, SLIDE_H - FOOTER_H, Inches(5), FOOTER_H, anchor=MSO_ANCHOR.MIDDLE)
    _run(lf.paragraphs[0], section, size=9, color=SLATE)
    _, rf = _txbox(slide, SLIDE_W - Inches(6) - MARGIN, SLIDE_H - FOOTER_H, Inches(6), FOOTER_H,
                   anchor=MSO_ANCHOR.MIDDLE)
    rf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _run(rf.paragraphs[0], f"Page {page} / {total}  ·  CONFIDENTIAL — BOARD ONLY", size=9, color=SLATE)

    if is_draft:
        _watermark(slide)


def _watermark(slide):
    wm = slide.shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(10.3), Inches(2.3))
    tf = wm.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = _run(p, "DRAFT", size=120, bold=True, color=RGBColor(0xE2, 0xE8, 0xF0))
    wm.rotation = -20


def _title(slide, title: str, key_message: str = ""):
    _, tf = _txbox(slide, MARGIN, Inches(0.62), SLIDE_W - 2 * MARGIN, Inches(0.85))
    _run(tf.paragraphs[0], title, size=24, bold=True, color=INK)
    if key_message:
        kp = tf.add_paragraph()
        _run(kp, key_message, size=13, color=SLATE)


# ---------------------------------------------------------------------------
# Table helper
# ---------------------------------------------------------------------------

def _table(slide, left, top, width, headers: List[str], rows: List[List[str]],
           col_w: Optional[List[float]] = None, rag_col: Optional[int] = None,
           rag_keys: Optional[List[str]] = None, row_h: float = 0.34):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = Inches(row_h * n_rows)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_w:
        total = sum(col_w)
        for ci, w in enumerate(col_w):
            gtable.columns[ci].width = Emu(int(width * (w / total)))

    # Header
    for ci, h in enumerate(headers):
        cell = gtable.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(4); cell.margin_right = Pt(4)
        cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
        para = cell.text_frame.paragraphs[0]
        _run(para, h, size=9.5, bold=True, color=WHITE)

    for ri, row in enumerate(rows, start=1):
        rag = (rag_keys[ri - 1] if rag_keys and ri - 1 < len(rag_keys) else None)
        for ci, val in enumerate(row):
            cell = gtable.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
            cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
            cell.fill.solid()
            if rag_col is not None and ci == rag_col and rag:
                cell.fill.fore_color.rgb = RAG_SOFT.get(rag, WHITE)
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xF8, 0xFA, 0xFC)
            para = cell.text_frame.paragraphs[0]
            color = RAG.get(rag, INK) if (rag_col is not None and ci == rag_col and rag) else INK
            _run(para, str(val), size=9.5, color=color,
                 bold=(rag_col is not None and ci == rag_col and bool(rag)))
    return gtable


# ---------------------------------------------------------------------------
# Per-slide renderers
# ---------------------------------------------------------------------------

def _bullets(slide, left, top, width, height, items, *, size=12, color=INK, bullet="•  "):
    _, tf = _txbox(slide, left, top, width, height)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        _run(p, bullet + str(it), size=size, color=color)


def _render_cover(slide, d):
    # Big centred title block (no standard title bar)
    _, tf = _txbox(slide, MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(2.4),
                   anchor=MSO_ANCHOR.MIDDLE)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tf.paragraphs[0], d.get("company_name", ""), size=40, bold=True, color=INK)
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
    _run(p, f"{d.get('report_type_label', 'Board Pack')}  ·  {d.get('period', '')}",
         size=20, color=SLATE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    _run(p2, d.get("sector", ""), size=14, color=MUTED)

    # IRR strip
    irr_base = d.get("irr_base"); irr_ic = d.get("irr_ic")
    if irr_base or irr_ic:
        _, itf = _txbox(slide, MARGIN, Inches(4.7), SLIDE_W - 2 * MARGIN, Inches(0.5))
        itf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(itf.paragraphs[0],
             f"Base (P50) IRR {irr_base or '—'}   vs   IC {irr_ic or '—'}"
             + (f"   ({d.get('irr_gap_bps'):+d} bps)" if d.get("irr_gap_bps") is not None else ""),
             size=14, bold=True, color=SLATE)

    _, ctf = _txbox(slide, MARGIN, Inches(5.5), SLIDE_W - 2 * MARGIN, Inches(0.9),
                    anchor=MSO_ANCHOR.MIDDLE)
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(ctf.paragraphs[0], d.get("classification", "Confidential — Board Only"),
         size=12, color=MUTED)
    p3 = ctf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    _run(p3, f"Prepared by {d.get('prepared_by', '')}  ·  {d.get('date', '')}",
         size=11, color=MUTED)


def _render_at_a_glance(slide, d):
    col_w = Inches(5.9)
    left2 = MARGIN + col_w + Inches(0.4)
    # KPI performance
    kpis = d.get("kpi_performance") or []
    _table(slide, MARGIN, BODY_TOP, col_w,
           ["KPI", "Actual", ""],
           [[k.get("label", ""), k.get("value", ""), ""] for k in kpis],
           col_w=[3, 2, 0.6], rag_col=2,
           rag_keys=[k.get("status") for k in kpis])

    # VCP milestone summary
    ms = d.get("vcp_milestones") or {}
    _, tf = _txbox(slide, MARGIN, BODY_TOP + Inches(2.2), col_w, Inches(2.0))
    _run(tf.paragraphs[0], "VCP Milestones", size=13, bold=True, color=NAVY)
    for label, key, rag in [("On track", "on_track", "green"), ("At risk", "at_risk", "amber"),
                            ("Behind", "behind", "red"), ("No data", "no_data", "grey")]:
        p = tf.add_paragraph()
        _run(p, f"{ms.get(key, 0)}  {label}", size=12, color=RAG.get(rag, INK), bold=True)

    # Peer position + IRR outlook
    bench = d.get("peer_position") or []
    _table(slide, left2, BODY_TOP, col_w,
           ["Peer metric", "Percentile"],
           [[b.get("metric", ""), b.get("percentile") or "—"] for b in bench],
           col_w=[3, 2])
    irr = d.get("irr_outlook") or {}
    _, itf = _txbox(slide, left2, BODY_TOP + Inches(2.2), col_w, Inches(1.6))
    _run(itf.paragraphs[0], "IRR Outlook", size=13, bold=True, color=NAVY)
    for lab, key in [("Bear", "bear"), ("Base", "base"), ("Bull", "bull"), ("IC", "ic")]:
        p = itf.add_paragraph()
        _run(p, f"{lab}: {irr.get(key) or '—'}", size=12, color=INK)

    if d.get("priority_action"):
        _, ptf = _txbox(slide, MARGIN, BODY_TOP + Inches(4.0), SLIDE_W - 2 * MARGIN, Inches(0.7))
        _run(ptf.paragraphs[0], "Priority action: ", size=12, bold=True, color=NAVY)
        _run(ptf.paragraphs[0], d.get("priority_action", ""), size=12, color=INK)


def _render_exec_summary(slide, d):
    _, tf = _txbox(slide, MARGIN, BODY_TOP, Inches(8.2), Inches(4.6))
    _run(tf.paragraphs[0], d.get("situation", ""), size=13, color=INK)
    if d.get("key_risks"):
        rp = tf.add_paragraph(); rp.space_before = Pt(8)
        _run(rp, "Key risks", size=12, bold=True, color=NAVY)
        for r in d.get("key_risks", []):
            p = tf.add_paragraph(); _run(p, "•  " + str(r), size=12, color=INK)
    if d.get("priority_action"):
        pp = tf.add_paragraph(); pp.space_before = Pt(8)
        _run(pp, "Priority action: ", size=12, bold=True, color=NAVY)
        _run(pp, d.get("priority_action", ""), size=12, color=INK)

    # Stat cards on the right
    stats = d.get("stats") or []
    top = BODY_TOP
    for s in stats:
        card = slide.shapes.add_shape(1, Inches(9.0), top, Inches(3.7), Inches(1.2))
        _fill(card, RGBColor(0xF1, 0xF5, 0xF9))
        _, ctf = _txbox(slide, Inches(9.15), top + Inches(0.08), Inches(3.4), Inches(1.05))
        _run(ctf.paragraphs[0], s.get("label", ""), size=9, bold=True, color=SLATE)
        vp = ctf.add_paragraph()
        _run(vp, s.get("value", ""), size=22, bold=True,
             color=RAG["red"] if s.get("negative") else INK)
        sp = ctf.add_paragraph()
        _run(sp, s.get("sub", ""), size=9, color=MUTED)
        top = top + Inches(1.4)


def _render_kpi_scorecard(slide, d):
    rows = d.get("rows") or []
    table_rows = [[r.get("metric", ""), r.get("actual", ""), r.get("ic_target", ""),
                   r.get("delta", ""), r.get("vs_last_qtr", ""), r.get("trend", "")]
                  for r in rows]
    _table(slide, MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN,
           ["Metric", "Actual", "IC Target", "Δ vs IC", "vs Last Qtr", "Trend"],
           table_rows, col_w=[3, 1.6, 1.6, 1.5, 1.6, 1.4],
           rag_col=0, rag_keys=[r.get("status") for r in rows])
    if d.get("legend"):
        _, lf = _txbox(slide, MARGIN, SLIDE_H - FOOTER_H - Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.4))
        _run(lf.paragraphs[0], d.get("legend", ""), size=9, italic=True, color=MUTED)


def _render_vcp_scorecard(slide, d):
    rows = d.get("rows") or []
    table_rows = [[r.get("initiative", ""), r.get("category", ""), r.get("baseline", ""),
                   r.get("target", ""), r.get("actual", ""), r.get("gap", ""),
                   r.get("status_label", "") + (" ⚠" if r.get("overdue") else ""), r.get("due", "")]
                  for r in rows]
    _table(slide, MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN,
           ["Initiative", "Category", "Base", "Target", "Actual", "Gap", "Status", "Due"],
           table_rows, col_w=[3.2, 1.6, 1.2, 1.2, 1.2, 1.2, 1.4, 1.0],
           rag_col=6, rag_keys=[r.get("status") for r in rows], row_h=0.32)


def _render_forward_curve(slide, d):
    png = _forward_curve_png(d)
    if png:
        slide.shapes.add_picture(io.BytesIO(png), MARGIN, BODY_TOP, width=Inches(8.4))
    # KPI strip on the right
    strip = d.get("kpi_strip") or []
    top = BODY_TOP
    for s in strip:
        _, tf = _txbox(slide, Inches(9.3), top, Inches(3.5), Inches(0.9))
        _run(tf.paragraphs[0], s.get("label", ""), size=9, bold=True, color=SLATE)
        vp = tf.add_paragraph(); _run(vp, s.get("value", ""), size=18, bold=True, color=INK)
        top = top + Inches(1.0)
    if d.get("methodology"):
        _, mf = _txbox(slide, MARGIN, SLIDE_H - FOOTER_H - Inches(0.55), Inches(8.4), Inches(0.5))
        _run(mf.paragraphs[0], d.get("methodology", ""), size=8, italic=True, color=MUTED)


def _render_benchmarks(slide, d):
    rows = d.get("rows") or []
    table_rows = [[r.get("metric", ""), r.get("company", ""), r.get("p25", ""),
                   r.get("median", ""), r.get("p75", ""), r.get("percentile") or "—",
                   r.get("gap", "")] for r in rows]
    _table(slide, MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN,
           ["Metric", "Company", "P25", "Median", "P75", "Percentile", "Gap"],
           table_rows, col_w=[3, 1.5, 1.3, 1.3, 1.3, 1.5, 1.3])
    sub = f"{d.get('sector_label', '')}"
    if d.get("peer_set_size"):
        sub += f"  ·  {d.get('peer_set_size')} peers"
    if d.get("estimated_band"):
        sub += "  ·  P25/P75 estimated"
    _, sf = _txbox(slide, MARGIN, SLIDE_H - FOOTER_H - Inches(0.5), SLIDE_W - 2 * MARGIN, Inches(0.4))
    _run(sf.paragraphs[0], sub, size=9, italic=True, color=MUTED)


def _render_risks(slide, d):
    if d.get("has_risks"):
        rows = d.get("rows") or []
        table_rows = [[str(r.get("priority", "")), r.get("status", ""), r.get("risk", ""),
                       (f"{r.get('irr_at_risk_bps'):+d}bps" if r.get("irr_at_risk_bps") is not None else "—"),
                       r.get("recommended_action", "")] for r in rows]
        rag_keys = [(r.get("status", "").lower() if r.get("status", "").lower() in RAG else "amber")
                    for r in rows]
        _table(slide, MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN,
               ["#", "Status", "Risk", "IRR at risk", "Recommended action"],
               table_rows, col_w=[0.5, 1.1, 4.2, 1.4, 4.5],
               rag_col=1, rag_keys=rag_keys, row_h=0.5)
    else:
        nb = d.get("no_risk_block") or {}
        _, tf = _txbox(slide, MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN, Inches(1.5))
        _run(tf.paragraphs[0], nb.get("headline", "No material risks this period"),
             size=16, bold=True, color=RAG["green"])
        if nb.get("detail"):
            p = tf.add_paragraph(); _run(p, nb.get("detail", ""), size=12, color=INK)

    bq = d.get("board_questions") or []
    if bq:
        _, qf = _txbox(slide, MARGIN, SLIDE_H - FOOTER_H - Inches(1.6), SLIDE_W - 2 * MARGIN, Inches(1.5))
        _run(qf.paragraphs[0], "Board questions", size=12, bold=True, color=NAVY)
        for q in bq[:3]:
            p = qf.add_paragraph(); _run(p, "•  " + str(q), size=10.5, color=SLATE)


def _render_irr_matrix(slide, d):
    png = _irr_heatmap_png(d.get("matrix") or {})
    if png:
        slide.shapes.add_picture(io.BytesIO(png), MARGIN, BODY_TOP, width=Inches(8.2))
    ic = d.get("ic_underwritten") or {}
    p50 = d.get("p50_trajectory") or {}
    _, tf = _txbox(slide, Inches(9.1), BODY_TOP, Inches(3.6), Inches(3.5))
    _run(tf.paragraphs[0], "IC underwritten", size=12, bold=True, color=NAVY)
    p = tf.add_paragraph()
    _run(p, f"{ic.get('exit_multiple', '')} · Yr {ic.get('year', '')} · {ic.get('irr', '')}",
         size=12, color=INK)
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _run(p2, "P50 trajectory", size=12, bold=True, color=NAVY)
    p3 = tf.add_paragraph()
    _run(p3, f"{p50.get('exit_multiple', '')} · Yr {p50.get('year', '')} · {p50.get('irr') or '—'}",
         size=12, color=INK)
    if p50.get("gap_bps") is not None:
        p4 = tf.add_paragraph()
        _run(p4, f"Gap: {p50.get('gap_bps'):+d} bps", size=12, bold=True,
             color=RAG["red"] if p50.get("gap_bps", 0) < 0 else RAG["green"])
    if d.get("legend"):
        lp = tf.add_paragraph(); lp.space_before = Pt(10)
        _run(lp, d.get("legend", ""), size=9, italic=True, color=MUTED)


def _render_audit(slide, d):
    prov = d.get("provenance") or {}
    _, tf = _txbox(slide, MARGIN, BODY_TOP, Inches(6.0), Inches(4.5))
    _run(tf.paragraphs[0], "Provenance", size=13, bold=True, color=NAVY)
    for label, key in [("Generated", "generated"), ("Approved by", "approved_by"),
                       ("Approved at", "approved_at"), ("Narrative", "narrative"),
                       ("Assembly", "assembly")]:
        if prov.get(key):
            p = tf.add_paragraph()
            _run(p, f"{label}: ", size=10.5, bold=True, color=SLATE)
            _run(p, str(prov.get(key)), size=10.5, color=INK)

    cits = d.get("citations") or []
    _, cf = _txbox(slide, Inches(7.0), BODY_TOP, Inches(5.7), Inches(4.5))
    _run(cf.paragraphs[0], "Citations & evidence", size=13, bold=True, color=NAVY)
    for c in cits:
        if isinstance(c, dict):
            label, source = c.get("label", ""), c.get("source", "")
            text = f"{label}: {source}" if label else source
        else:
            text = str(c)
        p = cf.add_paragraph(); _run(p, "•  " + text, size=10, color=INK)

    if d.get("disclaimer"):
        _, df = _txbox(slide, MARGIN, SLIDE_H - FOOTER_H - Inches(0.9), SLIDE_W - 2 * MARGIN, Inches(0.8))
        _run(df.paragraphs[0], d.get("disclaimer", ""), size=8, italic=True, color=MUTED)


_RENDERERS = {
    "Cover": _render_cover,
    "AtAGlance": _render_at_a_glance,
    "ExecSummary": _render_exec_summary,
    "KPIScorecard": _render_kpi_scorecard,
    "VCPScorecard": _render_vcp_scorecard,
    "ForwardCurve": _render_forward_curve,
    "Benchmarks": _render_benchmarks,
    "RisksActions": _render_risks,
    "IRRMatrix": _render_irr_matrix,
    "AuditTrail": _render_audit,
}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def generate_pptx(
    *,
    company_name: str,
    period: str,
    slides: List[Dict[str, Any]],
    is_approved: bool = False,
    notes: Optional[Dict[int, str]] = None,
) -> bytes:
    """Assemble the 10-slide editable Board Pack and return PPTX bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    is_draft = not is_approved
    notes = notes or {}
    total = len(slides)

    for slide_data in slides:
        slide = _add_blank(prs)
        stype = slide_data.get("type", "")
        data = slide_data.get("data", {}) or {}
        index = slide_data.get("index", 0)

        renderer = _RENDERERS.get(stype)
        if renderer:
            try:
                renderer(slide, data)
            except Exception as exc:  # never let one slide kill the export
                _, tf = _txbox(slide, MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN, Inches(1))
                _run(tf.paragraphs[0], f"[Slide render error: {exc}]", size=12, color=RAG["red"])

        # Title bar (Cover draws its own)
        if stype != "Cover":
            _title(slide, slide_data.get("title", ""), slide_data.get("key_message", ""))

        _chrome(slide, company=company_name, period=period, page=index + 1, total=total,
                section=slide_data.get("title", ""), is_draft=is_draft)

        # Speaker notes
        note_text = notes.get(index) or notes.get(str(index))
        if note_text:
            slide.notes_slide.notes_text_frame.text = note_text

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
